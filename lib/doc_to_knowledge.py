"""
PDF/PPT 批量转知识库工具
提取文档内容 + DeepSeek 智能总结

使用方法：
1. pip install openai pymupdf python-pptx
2. 设置环境变量 DEEPSEEK_API_KEY
3. 把这个脚本放到文档文件夹的父目录
4. 运行 python doc_to_knowledge.py

支持格式：PDF, PPT, PPTX
"""

import os
import sys
from pathlib import Path
from datetime import datetime
from openai import OpenAI

# PDF 处理
try:
    import fitz  # PyMuPDF
    PDF_SUPPORTED = True
except ImportError:
    PDF_SUPPORTED = False
    print("⚠️ 未安装 pymupdf，无法处理 PDF。运行: pip install pymupdf")

# PPT 处理
try:
    from pptx import Presentation
    PPT_SUPPORTED = True
except ImportError:
    PPT_SUPPORTED = False
    print("⚠️ 未安装 python-pptx，无法处理 PPT。运行: pip install python-pptx")

# 修复 Windows 中文编码
if sys.platform == 'win32':
    import locale
    locale.setlocale(locale.LC_ALL, '')

# ============ 配置 ============

DOC_EXTENSIONS = {'.pdf', '.ppt', '.pptx'}
OUTPUT_DIR = "txt_output"
DEEPSEEK_BASE_URL = "https://api.deepseek.com"
DEEPSEEK_MODEL = "deepseek-chat"
ENABLE_SUMMARY = True

# ============ 配置结束 ============


def get_script_dir():
    return Path(__file__).parent.resolve()


def find_all_docs(root_dir):
    """递归查找所有文档文件"""
    docs = []
    root_path = Path(root_dir)
    
    for item in root_path.iterdir():
        if item.name == OUTPUT_DIR:
            continue
            
        if item.is_dir():
            for doc_file in item.rglob('*'):
                if doc_file.suffix.lower() in DOC_EXTENSIONS:
                    docs.append(doc_file)
        elif item.is_file() and item.suffix.lower() in DOC_EXTENSIONS:
            docs.append(item)
    
    return sorted(docs)


def extract_pdf_text(pdf_path):
    """提取 PDF 文本内容"""
    if not PDF_SUPPORTED:
        return None, "PyMuPDF 未安装"
    
    try:
        doc = fitz.open(pdf_path)
        text_parts = []
        
        for page_num, page in enumerate(doc, 1):
            text = page.get_text()
            if text.strip():
                text_parts.append(f"=== 第 {page_num} 页 ===\n{text}")
        
        doc.close()
        
        full_text = "\n\n".join(text_parts)
        return full_text, None
        
    except Exception as e:
        return None, str(e)


def extract_ppt_text(ppt_path):
    """提取 PPT/PPTX 文本内容"""
    if not PPT_SUPPORTED:
        return None, "python-pptx 未安装"
    
    try:
        prs = Presentation(ppt_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            
            if slide_texts:
                text_parts.append(f"=== 幻灯片 {slide_num} ===\n" + "\n".join(slide_texts))
        
        full_text = "\n\n".join(text_parts)
        return full_text, None
        
    except Exception as e:
        return None, str(e)


def extract_text(doc_path):
    """根据文件类型提取文本"""
    suffix = doc_path.suffix.lower()
    
    if suffix == '.pdf':
        return extract_pdf_text(doc_path)
    elif suffix in ['.ppt', '.pptx']:
        return extract_ppt_text(doc_path)
    else:
        return None, f"不支持的格式: {suffix}"


def summarize_with_deepseek(client, text, doc_name, doc_type):
    """用 DeepSeek 总结提炼"""
    print(f"  [DeepSeek] 正在总结...")
    
    # 限制文本长度
    max_chars = 50000
    if len(text) > max_chars:
        text = text[:max_chars] + "\n\n[文本过长，已截取前半部分]"
    
    prompt = f"""你是一个专业的交易知识整理专家。请分析以下{doc_type}文档内容，提取核心知识点。

文档标题：{doc_name}

文档内容：
{text}

请按以下格式输出：

## 📌 核心主题
[一句话概括这个文档讲什么]

## 🎯 关键知识点
[列出 5-10 个最重要的知识点，每个用 1-2 句话说清楚]

## 📊 实战要点
[提取可以直接用于交易的具体方法、指标、规则]

## 💡 重要概念
[解释文档中提到的专业术语和概念]

## 🔗 知识关联
[这个内容和其他交易理论（如维科夫、K线、量价）的关系]

## 📋 要点摘要
[用 3-5 句话总结最核心的内容]

请用简洁专业的中文输出，重点突出可操作的交易知识。
"""

    try:
        response = client.chat.completions.create(
            model=DEEPSEEK_MODEL,
            messages=[
                {"role": "system", "content": "你是专业的交易知识整理专家，擅长从文档中提取核心交易知识。"},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            max_tokens=4000
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"  [DeepSeek] 总结失败: {e}")
        return None


def create_knowledge_card(doc_name, doc_type, summary, page_count=None):
    """生成知识卡片"""
    page_info = f"\n> 📄 页数：{page_count}" if page_count else ""
    
    card = f"""# {doc_name}

> 📅 生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}
> 📁 文档类型：{doc_type}{page_info}
> 🎬 来源：文档内容提取 + AI 总结

---

{summary}

---

*本文档由 AI 自动生成，用于知识库检索*
"""
    return card


def process_document(client, doc_path, output_dir, script_dir):
    """处理单个文档"""
    rel_path = doc_path.relative_to(script_dir)
    doc_name = doc_path.stem
    doc_type = "PDF" if doc_path.suffix.lower() == '.pdf' else "PPT"
    
    # 输出路径
    output_subdir = output_dir / rel_path.parent
    output_subdir.mkdir(parents=True, exist_ok=True)
    
    raw_path = output_subdir / f"{doc_name}_原始内容.txt"
    summary_path = output_subdir / f"{doc_name}_精华总结.txt"
    card_path = output_subdir / f"{doc_name}_知识卡片.md"
    
    # 检查是否已处理
    if card_path.exists():
        print(f"  ⏭️ 跳过（已存在）")
        return True
    
    try:
        # 1. 提取文本
        if raw_path.exists():
            print(f"  [提取] 读取已有内容...")
            with open(raw_path, "r", encoding="utf-8") as f:
                content = f.read()
            # 跳过头部信息
            if "\n\n" in content:
                text = content.split("\n\n", 1)[1]
            else:
                text = content
        else:
            print(f"  [提取] 正在提取文本...")
            text, error = extract_text(doc_path)
            
            if error:
                print(f"  ❌ 提取失败: {error}")
                return False
            
            if not text or not text.strip():
                print(f"  ⚠️ 文档内容为空或无法提取")
                return False
            
            # 保存原始内容
            with open(raw_path, "w", encoding="utf-8") as f:
                f.write(f"# {doc_name}\n")
                f.write(f"# 提取时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                f.write(f"# 源文件: {doc_path}\n\n")
                f.write(text)
            print(f"  ✅ 原始内容已保存")
        
        # 2. DeepSeek 总结
        if ENABLE_SUMMARY and client:
            summary = summarize_with_deepseek(client, text, doc_name, doc_type)
            
            if summary:
                # 保存精华总结
                with open(summary_path, "w", encoding="utf-8") as f:
                    f.write(summary)
                print(f"  ✅ 精华总结已保存")
                
                # 生成知识卡片
                card = create_knowledge_card(doc_name, doc_type, summary)
                with open(card_path, "w", encoding="utf-8") as f:
                    f.write(card)
                print(f"  ✅ 知识卡片已保存")
        else:
            # 没有 DeepSeek 时，直接生成简单卡片
            card = f"""# {doc_name}

> 📅 生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}
> 📁 文档类型：{doc_type}

---

{text[:5000]}

---
"""
            with open(card_path, "w", encoding="utf-8") as f:
                f.write(card)
            print(f"  ✅ 知识卡片已保存（无总结）")
        
        return True
        
    except Exception as e:
        print(f"  ❌ 处理失败: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    script_dir = get_script_dir()
    output_dir = script_dir / OUTPUT_DIR
    
    print("=" * 60)
    print("📚 PDF/PPT 批量转知识库工具")
    print("=" * 60)
    print(f"脚本目录: {script_dir}")
    print(f"输出目录: {output_dir}")
    print(f"支持格式: PDF, PPT, PPTX")
    print(f"DeepSeek 总结: {'启用' if ENABLE_SUMMARY else '禁用'}")
    print()
    
    # 检查依赖
    if not PDF_SUPPORTED:
        print("⚠️ PDF 支持未启用，运行: pip install pymupdf")
    if not PPT_SUPPORTED:
        print("⚠️ PPT 支持未启用，运行: pip install python-pptx")
    
    # 检查 DeepSeek API
    client = None
    if ENABLE_SUMMARY:
        api_key = os.environ.get("DEEPSEEK_API_KEY")
        if api_key:
            client = OpenAI(api_key=api_key, base_url=DEEPSEEK_BASE_URL)
            print("✅ DeepSeek API 已连接")
        else:
            print("⚠️ 未找到 DEEPSEEK_API_KEY，将只进行提取")
    
    # 查找文档
    print("\n正在扫描文档文件...")
    docs = find_all_docs(script_dir)
    
    if not docs:
        print("❌ 没有找到文档文件!")
        print(f"请确保 PDF/PPT 文件在 {script_dir} 的子目录中")
        return
    
    # 统计文件类型
    pdf_count = sum(1 for d in docs if d.suffix.lower() == '.pdf')
    ppt_count = sum(1 for d in docs if d.suffix.lower() in ['.ppt', '.pptx'])
    
    print(f"找到 {len(docs)} 个文档文件 (PDF: {pdf_count}, PPT: {ppt_count}):")
    for i, d in enumerate(docs[:10], 1):
        rel_path = d.relative_to(script_dir)
        print(f"  {i}. {rel_path}")
    if len(docs) > 10:
        print(f"  ... 还有 {len(docs) - 10} 个")
    
    print()
    input("按 Enter 开始处理...")
    
    # 处理文档
    success = 0
    failed = 0
    
    for i, doc_path in enumerate(docs, 1):
        print(f"\n[{i}/{len(docs)}] {doc_path.name}")
        
        if process_document(client, doc_path, output_dir, script_dir):
            success += 1
        else:
            failed += 1
    
    # 总结
    print("\n" + "=" * 60)
    print("🎉 处理完成!")
    print("=" * 60)
    print(f"成功: {success}")
    print(f"失败: {failed}")
    print(f"\n输出目录: {output_dir}")
    print(f"\n📌 下一步: 把 *_知识卡片.md 文件上传到 NotebookLM")


if __name__ == "__main__":
    main()
